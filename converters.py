import fitz  # PyMuPDF
import docx
import pandas as pd
import json
import xml.etree.ElementTree as ET
from PIL import Image
import os
import sys
import subprocess
from pathlib import Path
import tempfile
import shutil
import re
from bs4 import BeautifulSoup

# Professional conversion libraries
try:
    import pdf2docx
except ImportError:
    pdf2docx = None

try:
    import docx2pdf
except ImportError:
    docx2pdf = None

try:
    from weasyprint import HTML, CSS
except ImportError:
    HTML = CSS = None

try:
    import win32com.client
except ImportError:
    win32com = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import camelot
except ImportError:
    camelot = None

try:
    import tabula
except ImportError:
    tabula = None

try:
    import pikepdf
except ImportError:
    pikepdf = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

class FileConverter:
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()
        
    def __del__(self):
        try:
            shutil.rmtree(self.temp_dir, ignore_errors=True)
        except:
            pass

    def convert(self, input_path, output_format):
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input file not found: {input_path}")
            
        output_path = self._get_output_path(input_path, output_format)
        input_ext = Path(input_path).suffix.lower()
        output_ext = output_format.lower()
        
        try:
            # Ensure output directory exists
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            # PDF conversions
            if input_ext == '.pdf':
                if output_ext == 'docx':
                    return self._pdf_to_docx_professional(input_path, output_path)
                elif output_ext == 'html':
                    return self._pdf_to_html_professional(input_path, output_path)
                elif output_ext == 'txt':
                    return self._pdf_to_txt_professional(input_path, output_path)
                elif output_ext == 'xlsx':
                    return self._pdf_to_xlsx_professional(input_path, output_path)
                elif output_ext == 'csv':
                    return self._pdf_to_csv_professional(input_path, output_path)
                elif output_ext in ['jpg', 'jpeg', 'png']:
                    return self._pdf_to_image_professional(input_path, output_path)
                elif output_ext == 'pdfa':
                    return self._pdf_to_pdfa(input_path, output_path)
            
            # DOCX conversions
            elif input_ext == '.docx':
                if output_ext == 'pdf':
                    return self._docx_to_pdf_professional(input_path, output_path)
                elif output_ext == 'html':
                    return self._docx_to_html_professional(input_path, output_path)
                elif output_ext == 'txt':
                    return self._docx_to_txt_professional(input_path, output_path)
                elif output_ext == 'pptx':
                    return self._docx_to_pptx_with_images(input_path, output_path)
                else:
                    text = self._extract_docx_text(input_path)
                    return self._create_from_text(text, output_path, output_ext)
            
            # HTML conversions
            elif input_ext == '.html':
                if output_ext == 'pdf':
                    return self._html_to_pdf_professional(input_path, output_path)
                elif output_ext == 'docx':
                    return self._html_to_docx_professional(input_path, output_path)
                else:
                    text = self._extract_html_text(input_path)
                    return self._create_from_text(text, output_path, output_ext)
            
            # Excel conversions
            elif input_ext == '.xlsx':
                if output_ext == 'pdf':
                    return self._xlsx_to_pdf_professional(input_path, output_path)
                elif output_ext == 'html':
                    return self._xlsx_to_html_professional(input_path, output_path)
                elif output_ext == 'csv':
                    return self._xlsx_to_csv_professional(input_path, output_path)
                else:
                    df = self.load_dataframe(input_path)
                    return self._dataframe_to_format(df, output_path, output_ext)
            
            # CSV conversions
            elif input_ext == '.csv':
                if output_ext == 'xlsx':
                    return self._csv_to_xlsx_professional(input_path, output_path)
                elif output_ext == 'pdf':
                    return self._csv_to_pdf_professional(input_path, output_path)
                elif output_ext == 'html':
                    return self._csv_to_html_professional(input_path, output_path)
                else:
                    df = self.load_dataframe(input_path)
                    return self._dataframe_to_format(df, output_path, output_ext)
            
            # Image conversions
            elif input_ext in ['.jpg', '.jpeg', '.png']:
                if output_ext == 'pdf':
                    return self._image_to_pdf_professional(input_path, output_path)
                elif output_ext == 'docx':
                    return self._image_to_docx_professional(input_path, output_path)
                else:
                    return self._image_convert(input_path, output_path, output_ext)
            
            # PowerPoint conversions
            elif input_ext == '.pptx':
                if output_ext == 'pdf':
                    return self._pptx_to_pdf(input_path, output_path)
                else:
                    text = self._extract_pptx_text(input_path)
                    return self._create_from_text(text, output_path, output_ext)
            
            # Text conversions
            elif input_ext == '.txt':
                text = self._extract_txt_text(input_path)
                if output_ext == 'pdf':
                    return self._text_to_pdf(text, output_path)
                elif output_ext == 'pptx':
                    return self._text_to_pptx(text, output_path)
                else:
                    return self._create_from_text(text, output_path, output_ext)
            
            # JSON/XML conversions
            elif input_ext in ['.json', '.xml']:
                df = self.load_dataframe(input_path)
                return self._dataframe_to_format(df, output_path, output_ext)
            
            else:
                raise ValueError(f"Unsupported input format: {input_ext}")
            
            # Copy from temp directory to final destination if needed
            if output_path.startswith(self.temp_dir):
                try:
                    final_output = self._get_final_output_path(input_path, output_format)
                    shutil.copy2(output_path, final_output)
                    return final_output
                except Exception as copy_error:
                    print(f"Could not copy to final destination: {copy_error}")
                    return output_path
            
            return output_path
                
        except Exception as e:
            if os.path.exists(output_path):
                try:
                    os.remove(output_path)
                except:
                    pass
            raise e
    
    def _get_final_output_path(self, input_path, output_format):
        """Get the intended final output path"""
        input_file = Path(input_path)
        return str(input_file.parent / f"{input_file.stem}.{output_format}")

    def _get_output_path(self, input_path, output_format):
        input_file = Path(input_path)
        # Always use temp directory for faster processing
        temp_output = os.path.join(self.temp_dir, f"{input_file.stem}.{output_format}")
        return temp_output

    # Professional PDF conversions
    def _pdf_to_docx_professional(self, input_path, output_path):
        try:
            if pdf2docx:
                # Ensure output directory exists and is writable
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                
                # Check if file is already open/locked
                if os.path.exists(output_path):
                    try:
                        os.remove(output_path)
                    except PermissionError:
                        # Use temp file if original is locked
                        temp_output = output_path.replace('.docx', '_temp.docx')
                        output_path = temp_output
                
                from pdf2docx import Converter
                cv = Converter(input_path)
                cv.convert(output_path, start=0, end=None)
                cv.close()
                return output_path
        except Exception as e:
            print(f"pdf2docx failed: {e}")
        
        # Fallback
        return self._pdf_to_docx_fallback(input_path, output_path)

    def _pdf_to_html_professional(self, input_path, output_path):
        try:
            if pdfplumber:
                with pdfplumber.open(input_path) as pdf:
                    html_content = ['<!DOCTYPE html><html><head><meta charset="utf-8">']
                    html_content.append('<style>')
                    html_content.append('body { font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }')
                    html_content.append('table { border-collapse: collapse; width: 100%; margin: 20px 0; }')
                    html_content.append('td, th { border: 1px solid #ddd; padding: 8px; text-align: left; }')
                    html_content.append('th { background-color: #f2f2f2; font-weight: bold; }')
                    html_content.append('.page { page-break-after: always; margin-bottom: 40px; }')
                    html_content.append('</style></head><body>')
                    
                    for page_num, page in enumerate(pdf.pages):
                        html_content.append(f'<div class="page" id="page-{page_num + 1}">')
                        
                        # Extract tables first
                        tables = page.extract_tables()
                        if tables:
                            for table in tables:
                                if table and len(table) > 0:
                                    html_content.append('<table>')
                                    for i, row in enumerate(table):
                                        if row and any(cell for cell in row if cell):
                                            tag = 'th' if i == 0 else 'td'
                                            html_content.append('<tr>')
                                            for cell in row:
                                                html_content.append(f'<{tag}>{cell or ""}</{tag}>')
                                            html_content.append('</tr>')
                                    html_content.append('</table>')
                        
                        # Extract remaining text
                        text = page.extract_text()
                        if text:
                            paragraphs = text.split('\n\n')
                            for para in paragraphs:
                                if para.strip():
                                    html_content.append(f'<p>{para.strip()}</p>')
                        
                        html_content.append('</div>')
                    
                    html_content.append('</body></html>')
                    
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write('\n'.join(html_content))
                    return output_path
        except Exception as e:
            print(f"Professional PDF to HTML failed: {e}")
        
        return self._pdf_to_html_fallback(input_path, output_path)

    def _pdf_to_xlsx_professional(self, input_path, output_path):
        try:
            if pdfplumber:
                with pdfplumber.open(input_path) as pdf:
                    all_tables = []
                    
                    for page in pdf.pages:
                        tables = page.extract_tables()
                        for table in tables:
                            if table and len(table) > 0:
                                clean_table = []
                                for row in table:
                                    if row and any(cell for cell in row if cell):
                                        clean_row = [str(cell).strip() if cell else "" for cell in row]
                                        clean_table.append(clean_row)
                                if clean_table:
                                    all_tables.append(clean_table)
                    
                    if all_tables:
                        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                            for i, table in enumerate(all_tables):
                                if len(table) > 1:
                                    df = pd.DataFrame(table[1:], columns=table[0])
                                else:
                                    df = pd.DataFrame(table)
                                sheet_name = f'Table_{i+1}' if len(all_tables) > 1 else 'Sheet1'
                                df.to_excel(writer, sheet_name=sheet_name, index=False)
                        return output_path
        except Exception as e:
            print(f"Professional PDF to Excel failed: {e}")
        
        return self._pdf_to_xlsx_fallback(input_path, output_path)

    def _docx_to_pdf_professional(self, input_path, output_path):
        try:
            if docx2pdf and sys.platform == "win32":
                # Ensure output directory exists
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                
                # Remove existing file if present
                if os.path.exists(output_path):
                    try:
                        os.remove(output_path)
                    except PermissionError:
                        output_path = output_path.replace('.pdf', '_temp.pdf')
                
                docx2pdf.convert(input_path, output_path)
                return output_path
        except Exception as e:
            print(f"docx2pdf failed: {e}")
        
        try:
            if win32com and sys.platform == "win32":
                import pythoncom
                pythoncom.CoInitialize()
                
                word = None
                try:
                    word = win32com.client.Dispatch("Word.Application")
                    word.Visible = False
                    word.DisplayAlerts = False
                    
                    # Ensure paths are absolute
                    abs_input = os.path.abspath(input_path)
                    abs_output = os.path.abspath(output_path)
                    
                    # Remove existing output file
                    if os.path.exists(abs_output):
                        try:
                            os.remove(abs_output)
                        except PermissionError:
                            abs_output = abs_output.replace('.pdf', '_temp.pdf')
                    
                    doc = word.Documents.Open(abs_input, ReadOnly=True)
                    doc.SaveAs2(abs_output, FileFormat=17)  # Use SaveAs2
                    doc.Close()
                    
                    return abs_output
                    
                except Exception as com_error:
                    print(f"Word COM operation failed: {com_error}")
                    raise
                finally:
                    if word:
                        try:
                            word.Quit()
                        except:
                            pass
                    pythoncom.CoUninitialize()
                    
        except Exception as e:
            print(f"Word COM failed: {e}")
        
        return self._docx_to_pdf_fallback(input_path, output_path)

    def _html_to_pdf_professional(self, input_path, output_path):
        try:
            if HTML:
                HTML(filename=input_path).write_pdf(output_path)
                return output_path
        except Exception as e:
            print(f"WeasyPrint failed: {e}")
        
        return self._html_to_pdf_fallback(input_path, output_path)

    def _xlsx_to_pdf_professional(self, input_path, output_path):
        try:
            if win32com and sys.platform == "win32":
                import pythoncom
                pythoncom.CoInitialize()
                
                excel = None
                try:
                    excel = win32com.client.Dispatch("Excel.Application")
                    excel.Visible = False
                    excel.DisplayAlerts = False
                    
                    abs_input = os.path.abspath(input_path)
                    abs_output = os.path.abspath(output_path)
                    
                    # Remove existing output file
                    if os.path.exists(abs_output):
                        try:
                            os.remove(abs_output)
                        except PermissionError:
                            abs_output = abs_output.replace('.pdf', '_temp.pdf')
                    
                    wb = excel.Workbooks.Open(abs_input, ReadOnly=True)
                    wb.ExportAsFixedFormat(0, abs_output)  # Use ExportAsFixedFormat
                    wb.Close()
                    
                    return abs_output
                    
                except Exception as com_error:
                    print(f"Excel COM operation failed: {com_error}")
                    raise
                finally:
                    if excel:
                        try:
                            excel.Quit()
                        except:
                            pass
                    pythoncom.CoUninitialize()
                    
        except Exception as e:
            print(f"Excel COM failed: {e}")
        
        return self._xlsx_to_pdf_fallback(input_path, output_path)

    def _image_to_pdf_professional(self, input_path, output_path):
        try:
            img = Image.open(input_path)
            if img.mode in ('RGBA', 'LA', 'P'):
                img = img.convert('RGB')
            img.save(output_path, "PDF", resolution=300.0, quality=95)
            return output_path
        except Exception as e:
            print(f"Professional image to PDF failed: {e}")
            return self._image_to_pdf_fallback(input_path, output_path)

    # Fallback methods
    def _pdf_to_docx_fallback(self, input_path, output_path):
        doc_pdf = fitz.open(input_path)
        doc_docx = docx.Document()
        
        for page in doc_pdf:
            text = page.get_text()
            if text.strip():
                for para in text.split('\n\n'):
                    if para.strip():
                        doc_docx.add_paragraph(para.strip())
        
        doc_docx.save(output_path)
        doc_pdf.close()
        return output_path

    def _pdf_to_html_fallback(self, input_path, output_path):
        doc = fitz.open(input_path)
        html_content = ['<html><head><meta charset="utf-8"></head><body>']
        
        for page_num, page in enumerate(doc):
            text = page.get_text()
            html_content.append(f'<div class="page-{page_num + 1}">')
            for para in text.split('\n\n'):
                if para.strip():
                    html_content.append(f'<p>{para.strip()}</p>')
            html_content.append('</div>')
        
        html_content.append('</body></html>')
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(html_content))
        doc.close()
        return output_path

    def _pdf_to_xlsx_fallback(self, input_path, output_path):
        doc = fitz.open(input_path)
        text_lines = []
        
        for page in doc:
            text_lines.extend([line.strip() for line in page.get_text().split('\n') if line.strip()])
        
        df = pd.DataFrame(text_lines, columns=['Content'])
        df.to_excel(output_path, index=False)
        doc.close()
        return output_path

    def _docx_to_pdf_fallback(self, input_path, output_path):
        doc = docx.Document(input_path)
        
        pdf_doc = fitz.open()
        page = pdf_doc.new_page()
        
        y_position = 50
        for para in doc.paragraphs:
            if para.text.strip():
                if y_position > 750:
                    page = pdf_doc.new_page()
                    y_position = 50
                page.insert_text((50, y_position), para.text, fontsize=11)
                y_position += 20
        
        pdf_doc.save(output_path)
        pdf_doc.close()
        return output_path

    def _html_to_pdf_fallback(self, input_path, output_path):
        with open(input_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        text = soup.get_text()
        
        return self._text_to_pdf(text, output_path)

    def _xlsx_to_pdf_fallback(self, input_path, output_path):
        df = pd.read_excel(input_path)
        
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
        from reportlab.lib import colors
        
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        elements = []
        
        data = [df.columns.tolist()] + df.values.tolist()
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        elements.append(table)
        doc.build(elements)
        return output_path

    def _image_to_pdf_fallback(self, input_path, output_path):
        img = Image.open(input_path)
        if img.mode != 'RGB':
            img = img.convert('RGB')
        img.save(output_path, "PDF")
        return output_path

    # Text extraction methods
    def extract_text(self, file_path):
        ext = Path(file_path).suffix.lower()
        if ext == '.pdf':
            return self._extract_pdf_text(file_path)
        elif ext == '.docx':
            return self._extract_docx_text(file_path)
        elif ext == '.txt':
            return self._extract_txt_text(file_path)
        elif ext == '.html':
            return self._extract_html_text(file_path)
        else:
            return ""

    def _extract_pdf_text(self, file_path):
        doc = fitz.open(file_path)
        text_parts = []
        for page in doc:
            # Use layout preservation for better spacing
            blocks = page.get_text("dict")["blocks"]
            page_lines = []
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        line_text = ""
                        for span in line["spans"]:
                            line_text += span["text"]
                        if line_text.strip():
                            page_lines.append(line_text.rstrip())
            if page_lines:
                text_parts.append('\n'.join(page_lines))
        doc.close()
        return '\n\n'.join(text_parts)

    def _extract_docx_text(self, file_path):
        doc = docx.Document(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                # Normalize spacing while preserving structure
                clean_text = ' '.join(para.text.split())
                text_parts.append(clean_text)
        
        for table in doc.tables:
            for row in table.rows:
                row_cells = [' '.join(cell.text.split()) for cell in row.cells]
                row_text = '\t'.join(row_cells)
                if row_text.strip():
                    text_parts.append(row_text)
        
        return '\n'.join(text_parts)

    def _extract_txt_text(self, file_path):
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise ValueError("Could not decode text file")

    def _extract_html_text(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        soup = BeautifulSoup(content, 'html.parser')
        for script in soup(["script", "style"]):
            script.decompose()
        
        return soup.get_text()

    def _extract_pptx_text(self, file_path):
        if not Presentation:
            raise ImportError("python-pptx required")
        
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return '\n'.join(text_parts)

    # Data loading methods
    def load_dataframe(self, file_path):
        ext = Path(file_path).suffix.lower()
        
        if ext == '.csv':
            return pd.read_csv(file_path)
        elif ext == '.xlsx':
            return pd.read_excel(file_path)
        elif ext == '.json':
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            if isinstance(data, list):
                return pd.DataFrame(data)
            else:
                return pd.DataFrame([data])
        elif ext == '.xml':
            tree = ET.parse(file_path)
            root = tree.getroot()
            rows = []
            for item in root:
                row = {child.tag: child.text for child in item}
                rows.append(row)
            return pd.DataFrame(rows)
        else:
            raise ValueError(f"Unsupported data format: {ext}")

    # Output creation methods
    def _create_from_text(self, text, output_path, ext):
        if ext == 'txt':
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
        elif ext == 'docx':
            doc = docx.Document()
            for line in text.split('\n'):
                if line.strip():
                    doc.add_paragraph(line)
            doc.save(output_path)
        elif ext == 'html':
            html = f'<html><head><meta charset="utf-8"></head><body><pre>{text}</pre></body></html>'
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html)
        elif ext == 'csv':
            lines = [line for line in text.split('\n') if line.strip()]
            df = pd.DataFrame(lines, columns=['Content'])
            df.to_csv(output_path, index=False)
        elif ext == 'xlsx':
            lines = [line for line in text.split('\n') if line.strip()]
            df = pd.DataFrame(lines, columns=['Content'])
            df.to_excel(output_path, index=False)
        elif ext == 'json':
            lines = [line for line in text.split('\n') if line.strip()]
            data = {'content': lines}
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
        elif ext == 'xml':
            lines = [line for line in text.split('\n') if line.strip()]
            root = ET.Element("Document")
            for line in lines:
                ET.SubElement(root, "Line").text = line
            tree = ET.ElementTree(root)
            tree.write(output_path, encoding='utf-8', xml_declaration=True)
        elif ext == 'pptx':
            return self._text_to_pptx(text, output_path)
        
        return output_path

    def _dataframe_to_format(self, df, output_path, ext):
        if ext == 'csv':
            df.to_csv(output_path, index=False)
        elif ext == 'xlsx':
            df.to_excel(output_path, index=False)
        elif ext == 'json':
            df.to_json(output_path, orient='records', indent=2, force_ascii=False)
        elif ext == 'xml':
            root = ET.Element("Root")
            for _, row in df.iterrows():
                item = ET.SubElement(root, "Row")
                for col, val in row.items():
                    ET.SubElement(item, str(col)).text = str(val) if pd.notna(val) else ""
            tree = ET.ElementTree(root)
            tree.write(output_path, encoding='utf-8', xml_declaration=True)
        elif ext == 'pdf':
            return self._dataframe_to_pdf(df, output_path)
        elif ext == 'pptx':
            text = df.to_string(index=False)
            return self._text_to_pptx(text, output_path)
        else:
            text = df.to_string(index=False)
            return self._create_from_text(text, output_path, ext)
        
        return output_path

    def _dataframe_to_pdf(self, df, output_path):
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
        from reportlab.lib import colors
        
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        elements = []
        
        data = [df.columns.tolist()] + df.values.tolist()
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        elements.append(table)
        doc.build(elements)
        return output_path

    def _text_to_pdf(self, text, output_path):
        doc = fitz.open()
        page = doc.new_page()
        
        y_position = 50
        for line in text.split('\n'):
            if y_position > 750:
                page = doc.new_page()
                y_position = 50
            page.insert_text((50, y_position), line, fontsize=11)
            y_position += 15
        
        doc.save(output_path)
        doc.close()
        return output_path

    # Additional professional methods
    def _pdf_to_txt_professional(self, input_path, output_path):
        try:
            if pdfplumber:
                with pdfplumber.open(input_path) as pdf:
                    text_parts = []
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            text_parts.append(text)
                    
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write('\n\n'.join(text_parts))
                    return output_path
        except:
            pass
        
        text = self._extract_pdf_text(input_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return output_path

    def _pdf_to_csv_professional(self, input_path, output_path):
        try:
            if pdfplumber:
                with pdfplumber.open(input_path) as pdf:
                    all_rows = []
                    for page in pdf.pages:
                        tables = page.extract_tables()
                        for table in tables:
                            if table:
                                all_rows.extend(table)
                    
                    if all_rows and len(all_rows) > 1:
                        df = pd.DataFrame(all_rows[1:], columns=all_rows[0])
                        df.to_csv(output_path, index=False)
                        return output_path
        except:
            pass
        
        text = self._extract_pdf_text(input_path)
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        df = pd.DataFrame(lines, columns=['Content'])
        df.to_csv(output_path, index=False)
        return output_path

    def _pdf_to_image_professional(self, input_path, output_path):
        doc = fitz.open(input_path)
        page = doc[0]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # High resolution
        pix.save(output_path)
        doc.close()
        return output_path

    def _docx_to_html_professional(self, input_path, output_path):
        doc = docx.Document(input_path)
        html_content = ['<html><head><meta charset="utf-8"></head><body>']
        
        for para in doc.paragraphs:
            if para.text.strip():
                html_content.append(f'<p>{para.text}</p>')
        
        for table in doc.tables:
            html_content.append('<table border="1">')
            for row in table.rows:
                html_content.append('<tr>')
                for cell in row.cells:
                    html_content.append(f'<td>{cell.text}</td>')
                html_content.append('</tr>')
            html_content.append('</table>')
        
        html_content.append('</body></html>')
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(html_content))
        return output_path

    def _docx_to_txt_professional(self, input_path, output_path):
        text = self._extract_docx_text(input_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return output_path

    def _html_to_docx_professional(self, input_path, output_path):
        with open(input_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        doc = docx.Document()
        
        for element in soup.find_all(['p', 'h1', 'h2', 'h3', 'div']):
            text = element.get_text().strip()
            if text:
                doc.add_paragraph(text)
        
        doc.save(output_path)
        return output_path

    def _csv_to_xlsx_professional(self, input_path, output_path):
        df = pd.read_csv(input_path)
        df.to_excel(output_path, index=False)
        return output_path

    def _csv_to_pdf_professional(self, input_path, output_path):
        df = pd.read_csv(input_path)
        return self._dataframe_to_pdf(df, output_path)

    def _csv_to_html_professional(self, input_path, output_path):
        df = pd.read_csv(input_path)
        html = df.to_html(index=False, table_id='data-table')
        
        full_html = f"""
        <html>
        <head>
            <meta charset="utf-8">
            <style>
                body {{ font-family: Arial, sans-serif; margin: 40px; }}
                table {{ border-collapse: collapse; width: 100%; }}
                th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
                th {{ background-color: #f2f2f2; }}
            </style>
        </head>
        <body>
            {html}
        </body>
        </html>
        """
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_html)
        return output_path

    def _xlsx_to_csv_professional(self, input_path, output_path):
        df = pd.read_excel(input_path)
        df.to_csv(output_path, index=False)
        return output_path

    def _xlsx_to_html_professional(self, input_path, output_path):
        df = pd.read_excel(input_path)
        html = df.to_html(index=False, table_id='data-table')
        
        full_html = f"""
        <html>
        <head>
            <meta charset="utf-8">
            <style>
                body {{ font-family: Arial, sans-serif; margin: 40px; }}
                table {{ border-collapse: collapse; width: 100%; }}
                th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
                th {{ background-color: #f2f2f2; }}
            </style>
        </head>
        <body>
            {html}
        </body>
        </html>
        """
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_html)
        return output_path

    def _image_to_docx_professional(self, input_path, output_path):
        doc = docx.Document()
        
        # Try OCR first for text extraction with formatting
        if pytesseract:
            try:
                img = Image.open(input_path)
                # Use TSV output to preserve positioning
                ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT, config='--psm 6')
                
                # Group text by lines based on top coordinate
                lines = {}
                for i in range(len(ocr_data['text'])):
                    if int(ocr_data['conf'][i]) > 30:  # Confidence threshold
                        text = ocr_data['text'][i].strip()
                        if text:
                            top = ocr_data['top'][i]
                            left = ocr_data['left'][i]
                            
                            # Group by line (similar top values)
                            line_key = top // 10 * 10  # Group within 10 pixels
                            if line_key not in lines:
                                lines[line_key] = []
                            lines[line_key].append((left, text))
                
                # Sort lines by vertical position and create formatted text
                for line_top in sorted(lines.keys()):
                    # Sort words in line by horizontal position
                    line_words = sorted(lines[line_top], key=lambda x: x[0])
                    
                    # Calculate spacing between words
                    line_text = ""
                    prev_right = 0
                    for left, text in line_words:
                        # Add spacing based on gap
                        if prev_right > 0:
                            gap = left - prev_right
                            if gap > 50:  # Large gap = tab
                                line_text += "\t"
                            elif gap > 20:  # Medium gap = multiple spaces
                                line_text += "  "
                            else:
                                line_text += " "
                        line_text += text
                        prev_right = left + len(text) * 8  # Approximate character width
                    
                    if line_text.strip():
                        doc.add_paragraph(line_text)
                
                # Add original image as well
                doc.add_paragraph("\n[Original Image]")
                doc.add_picture(input_path, width=docx.shared.Inches(6))
                
            except Exception as e:
                print(f"OCR failed: {e}")
                # Fallback to image only
                doc.add_picture(input_path, width=docx.shared.Inches(6))
        else:
            # No OCR available, just embed image
            doc.add_picture(input_path, width=docx.shared.Inches(6))
        
        doc.save(output_path)
        return output_path

    def _image_convert(self, input_path, output_path, output_ext):
        # For text formats, try OCR with formatting preservation
        if output_ext in ['txt', 'docx', 'html'] and pytesseract:
            try:
                img = Image.open(input_path)
                
                if output_ext == 'txt':
                    # Use data extraction for better spacing
                    ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT, 
                                                        config='--psm 6 -c preserve_interword_spaces=1')
                    
                    # Reconstruct text with proper spacing
                    lines = {}
                    for i in range(len(ocr_data['text'])):
                        if int(ocr_data['conf'][i]) > 30:
                            text = ocr_data['text'][i].strip()
                            if text:
                                top = ocr_data['top'][i]
                                left = ocr_data['left'][i]
                                line_key = top // 5 * 5  # Group by line
                                if line_key not in lines:
                                    lines[line_key] = []
                                lines[line_key].append((left, text))
                    
                    # Build text with preserved spacing
                    formatted_text = []
                    for line_top in sorted(lines.keys()):
                        line_words = sorted(lines[line_top], key=lambda x: x[0])
                        line_text = ""
                        prev_right = 0
                        
                        for left, word in line_words:
                            if prev_right > 0:
                                gap = left - prev_right
                                if gap > 50:  # Large gap
                                    line_text += "\t"
                                elif gap > 15:  # Medium gap
                                    line_text += "  "
                                else:  # Small gap
                                    line_text += " "
                            line_text += word
                            prev_right = left + len(word) * 8
                        
                        if line_text.strip():
                            formatted_text.append(line_text)
                    
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write('\n'.join(formatted_text))
                    return output_path
                    
                elif output_ext == 'docx':
                    return self._image_to_docx_professional(input_path, output_path)
                    
                elif output_ext == 'html':
                    # Extract with positioning data
                    ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT, config='--psm 6')
                    
                    html_content = ['<html><head><meta charset="utf-8"><style>']  
                    html_content.append('.ocr-line { margin: 2px 0; white-space: pre; font-family: monospace; }')
                    html_content.append('</style></head><body>')
                    
                    # Group by lines
                    lines = {}
                    for i in range(len(ocr_data['text'])):
                        if int(ocr_data['conf'][i]) > 30:
                            text = ocr_data['text'][i].strip()
                            if text:
                                top = ocr_data['top'][i]
                                left = ocr_data['left'][i]
                                line_key = top // 10 * 10
                                if line_key not in lines:
                                    lines[line_key] = []
                                lines[line_key].append((left, text))
                    
                    # Create HTML with preserved spacing
                    for line_top in sorted(lines.keys()):
                        line_words = sorted(lines[line_top], key=lambda x: x[0])
                        line_text = ""
                        prev_right = 0
                        for left, text in line_words:
                            if prev_right > 0:
                                gap = left - prev_right
                                if gap > 50:
                                    line_text += "\t"
                                elif gap > 20:
                                    line_text += "  "
                                else:
                                    line_text += " "
                            line_text += text
                            prev_right = left + len(text) * 8
                        
                        if line_text.strip():
                            html_content.append(f'<div class="ocr-line">{line_text}</div>')
                    
                    html_content.append('</body></html>')
                    
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write('\n'.join(html_content))
                    return output_path
                    
            except Exception as e:
                print(f"OCR conversion failed: {e}")
        
        # Fallback to regular image conversion
        img = Image.open(input_path)
        
        if output_ext in ['jpg', 'jpeg'] and img.mode in ('RGBA', 'LA'):
            img = img.convert('RGB')
        
        img.save(output_path)
        return output_path

    def _pptx_to_pdf(self, input_path, output_path):
        text = self._extract_pptx_text(input_path)
        return self._text_to_pdf(text, output_path)

    def _pdf_to_pdfa(self, input_path, output_path):
        if pikepdf:
            with pikepdf.open(input_path) as pdf:
                pdf.save(output_path, linearize=True)
            return output_path
        else:
            # Simple copy if pikepdf not available
            shutil.copy2(input_path, output_path)
            return output_path
    
    def _text_to_pptx(self, text, output_path):
        if not Presentation:
            raise ImportError("python-pptx required for PPTX creation")
        
        prs = Presentation()
        
        # Split text into slides (by double newlines or every 10 lines)
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        if not paragraphs:
            paragraphs = [line.strip() for line in text.split('\n') if line.strip()]
        
        # Group paragraphs into slides (max 10 lines per slide)
        slides_content = []
        current_slide = []
        
        for para in paragraphs:
            current_slide.append(para)
            if len(current_slide) >= 10:
                slides_content.append(current_slide)
                current_slide = []
        
        if current_slide:
            slides_content.append(current_slide)
        
        # Create slides
        for i, slide_lines in enumerate(slides_content):
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = f"Slide {i + 1}"
            
            # Set content
            content = slide.placeholders[1]
            content.text = '\n'.join(slide_lines)
        
        prs.save(output_path)
        return output_path
    
    def _docx_to_pptx_with_images(self, input_path, output_path):
        if not Presentation:
            raise ImportError("python-pptx required for PPTX creation")
        
        doc = docx.Document(input_path)
        prs = Presentation()
        
        # Extract text and images from DOCX
        slide_content = []
        current_slide_text = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                current_slide_text.append(para.text.strip())
                # Create new slide every 10 paragraphs
                if len(current_slide_text) >= 10:
                    slide_content.append(current_slide_text)
                    current_slide_text = []
        
        if current_slide_text:
            slide_content.append(current_slide_text)
        
        # Extract images from DOCX
        images = []
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    image_path = os.path.join(self.temp_dir, f"img_{len(images)}.png")
                    with open(image_path, 'wb') as f:
                        f.write(image_data)
                    images.append(image_path)
                except:
                    pass
        
        # Create slides with text
        for i, slide_lines in enumerate(slide_content):
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = f"Slide {i + 1}"
            
            content = slide.placeholders[1]
            content.text = '\n'.join(slide_lines)
        
        # Add image slides
        for img_path in images:
            try:
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add image to slide
                left = prs.slide_width // 4
                top = prs.slide_height // 4
                width = prs.slide_width // 2
                height = prs.slide_height // 2
                
                slide.shapes.add_picture(img_path, left, top, width, height)
            except:
                pass
        
        prs.save(output_path)
        
        # Clean up temp images
        for img_path in images:
            try:
                os.remove(img_path)
            except:
                pass
        
        return output_path